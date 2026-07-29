"""
Tamil Nadu Elections 2026 — Dashboard Builder
Generates:
  1. 6 high-quality chart PNGs for Power BI / PPT
  2. Enriched CSV data model (7 files) for Power BI
  3. Professional 10-slide PowerPoint presentation

Run: python build_dashboard.py
"""

import os, sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
import numpy as np
import pandas as pd
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ═══════════════════════════════════════════════════════════════════════
# CONFIGURATION
# ═══════════════════════════════════════════════════════════════════════
PROJECT_DIR = os.path.dirname(os.path.abspath(__file__))
BASE = os.path.join(PROJECT_DIR, "data")
OUT  = os.path.join(BASE, "outputs")
os.makedirs(OUT, exist_ok=True)

PARTY_COLORS = {
    'DMK':    '#E8192C',
    'AIADMK': '#1B6B35',
    'INC':    '#0066CC',
    'TVK':    '#FF6B00',
    'PMK':    '#8B0000',
    'VCK':    '#4B0082',
    'BJP':    '#FF9933',
    'Others': '#9E9E9E',
}
TOP_PARTIES = ['TVK', 'DMK', 'AIADMK', 'INC', 'PMK', 'VCK', 'BJP']

# Report identity — deep navy + muted gold, an editorial/annual-report register
# rather than a generic bright-indigo "startup dashboard" look.
NAVY = '#12213F'
INK  = '#2B2B2B'
SOURCE_NOTE = 'Source: Tamil Nadu Election Commission (ECI)  ·  TN Election 2026 Analysis'

plt.rcParams.update({
    'font.family'      : 'Segoe UI',
    'font.size'        : 11,
    'text.color'       : INK,
    'axes.spines.top'  : False,
    'axes.spines.right': False,
    'axes.spines.left' : False,
    'axes.edgecolor'   : '#B5B5B5',
    'figure.facecolor' : 'white',
    'axes.facecolor'   : 'white',
    'axes.grid'        : True,
    'axes.axisbelow'   : True,
    'axes.grid.axis'   : 'y',
    'grid.alpha'       : 0.6,
    'grid.color'       : '#E2E2E2',
    'grid.linewidth'   : 0.7,
})
DPI = 200

# PPT colour palette
C_DARK   = RGBColor(0x12, 0x21, 0x3F)   # deep navy — report identity
C_ORANGE = RGBColor(0xA9, 0x81, 0x2F)   # muted gold accent (was bright orange)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT   = RGBColor(0x21, 0x21, 0x21)
C_GRAY   = RGBColor(0x75, 0x75, 0x75)
C_LIGHT  = RGBColor(0xF1, 0xF2, 0xF6)
C_GREEN  = RGBColor(0x2E, 0x7D, 0x32)
C_RED    = RGBColor(0xC6, 0x28, 0x28)

# ═══════════════════════════════════════════════════════════════════════
# 1. LOAD & ENRICH
# ═══════════════════════════════════════════════════════════════════════
print("Loading data...")
df21   = pd.read_csv(os.path.join(BASE, 'tn_2021_results.csv'))
df26   = pd.read_csv(os.path.join(BASE, 'tn_2026_results.csv'))
master = pd.read_csv(os.path.join(BASE, 'constituency_master.csv'))

for df in [df21, df26, master]:
    df.columns = df.columns.str.strip().str.lower()

df21 = df21[df21['party'].str.upper() != 'NOTA'].copy()
df26 = df26[df26['party'].str.upper() != 'NOTA'].copy()

def enrich(df, year):
    df = df.copy()
    total = df.groupby('ac_number')['votes'].sum()
    df['total_votes'] = df['ac_number'].map(total)
    df['vote_share']  = (df['votes'] / df['total_votes'] * 100).round(2)
    df['rank']        = df.groupby('ac_number')['votes'].rank(ascending=False, method='min').astype(int)
    df['is_winner']   = (df['rank'] == 1).astype(int)
    df['year']        = year
    df['party_group'] = df['party'].apply(lambda p: p if p in TOP_PARTIES else 'Others')
    return df

df21 = enrich(df21, 2021)
df26 = enrich(df26, 2026)

# Add district from master (only column not already in results)
df21 = df21.merge(master[['ac_number', 'district']], on='ac_number', how='left')
df26 = df26.merge(master[['ac_number', 'district']], on='ac_number', how='left')

# Extract winners (idxmax ensures one winner per constituency even if tied)
win21 = df21.loc[df21.groupby('ac_number')['votes'].idxmax()].copy().reset_index(drop=True)
win26 = df26.loc[df26.groupby('ac_number')['votes'].idxmax()].copy().reset_index(drop=True)
win21['region'] = win21['region'].fillna('Unknown')
win26['region'] = win26['region'].fillna('Unknown')

# ── Side-by-side comparison ──────────────────────────────────────────
COLS21 = ['ac_number', 'constituency', 'district', 'region', 'reserved',
          'candidate', 'party', 'party_group', 'votes', 'total_votes', 'vote_share']
COLS26 = ['ac_number', 'candidate', 'party', 'party_group', 'votes', 'total_votes', 'vote_share']

cmp = (
    win21[COLS21].rename(columns={
        'candidate': 'candidate_21', 'party': 'party_21', 'party_group': 'party_21_group',
        'votes': 'votes_21', 'total_votes': 'total_21', 'vote_share': 'share_21'
    }).merge(
        win26[COLS26].rename(columns={
            'candidate': 'candidate_26', 'party': 'party_26', 'party_group': 'party_26_group',
            'votes': 'votes_26', 'total_votes': 'total_26', 'vote_share': 'share_26'
        }), on='ac_number'
    )
)
cmp['flipped']     = (cmp['party_21'] != cmp['party_26']).astype(int)
cmp['share_swing'] = (cmp['share_26'] - cmp['share_21']).round(2)

# ── Key metrics ──────────────────────────────────────────────────────
n_const    = len(cmp)
n_flipped  = int(cmp['flipped'].sum())
avg21      = round(win21['vote_share'].mean(), 1)
avg26      = round(win26['vote_share'].mean(), 1)
above50_21 = int((win21['vote_share'] > 50).sum())
above50_26 = int((win26['vote_share'] > 50).sum())
below35_21 = int((win21['vote_share'] < 35).sum())
below35_26 = int((win26['vote_share'] < 35).sum())
party_seats26 = win26.groupby('party').size()
top_party  = party_seats26.idxmax()
top_seats  = int(party_seats26.max())

print(f"✅ {n_const} constituencies | {n_flipped} flipped | Top: {top_party} ({top_seats} seats)")
print(f"   Avg share: 2021={avg21}%  2026={avg26}%")

# ═══════════════════════════════════════════════════════════════════════
# 2. POWER BI DATA MODEL EXPORTS
# ═══════════════════════════════════════════════════════════════════════
print("\nExporting Power BI data model...")

# Fact: all candidate votes (both years combined)
pd.concat([df21, df26], ignore_index=True).to_csv(
    os.path.join(OUT, 'fact_all_votes.csv'), index=False)

# Fact: winners side-by-side comparison
cmp.to_csv(os.path.join(OUT, 'fact_winners_comparison.csv'), index=False)

# Fact: regional summary per year
reg_rows = []
for yr, w in [(2021, win21), (2026, win26)]:
    g = w.groupby(['region', 'party_group']).agg(
        seats=('ac_number', 'count'),
        avg_share=('vote_share', 'mean'),
        total_votes=('votes', 'sum')
    ).reset_index()
    g['year'] = yr
    reg_rows.append(g)
reg_summary = pd.concat(reg_rows, ignore_index=True)
reg_summary['avg_share'] = reg_summary['avg_share'].round(2)
reg_summary.to_csv(os.path.join(OUT, 'fact_regional_summary.csv'), index=False)

# Fact: party summary per year
party_rows = []
for yr, w in [(2021, win21), (2026, win26)]:
    g = w.groupby('party_group').agg(
        seats=('ac_number', 'count'),
        avg_share=('vote_share', 'mean'),
        total_votes=('votes', 'sum')
    ).reset_index()
    g['year'] = yr
    party_rows.append(g)
party_summary = pd.concat(party_rows, ignore_index=True)
party_summary['avg_share'] = party_summary['avg_share'].round(2)
party_summary.to_csv(os.path.join(OUT, 'fact_party_summary.csv'), index=False)

# Fact: flip matrix
flip_matrix = (
    cmp[cmp['flipped'] == 1]
    .groupby(['party_21_group', 'party_26_group'])
    .size()
    .reset_index(name='seats_flipped')
)
flip_matrix.to_csv(os.path.join(OUT, 'fact_flip_matrix.csv'), index=False)

# Dim: constituency
master.to_csv(os.path.join(OUT, 'dim_constituency.csv'), index=False)

# Dim: party metadata
pd.DataFrame({
    'party': TOP_PARTIES,
    'color': [PARTY_COLORS[p] for p in TOP_PARTIES],
    'alliance_2021': ['TVK Alliance','DMK Alliance','AIADMK Alliance','DMK Alliance',
                      'AIADMK Alliance','DMK Alliance','AIADMK Alliance'],
    'alliance_2026': ['TVK Alliance','DMK Alliance','AIADMK Alliance','DMK Alliance',
                      'AIADMK Alliance','DMK Alliance','AIADMK Alliance'],
}).to_csv(os.path.join(OUT, 'dim_party.csv'), index=False)

print("✅ Power BI data model: 7 files exported")

# ═══════════════════════════════════════════════════════════════════════
# 3. CHART GENERATION
# ═══════════════════════════════════════════════════════════════════════
print("\nGenerating charts...")

def save_chart(fig, name, source=SOURCE_NOTE):
    fig.text(0.01, 0.01, source, fontsize=7.5, color='#9A9A9A',
              style='italic', ha='left')
    path = os.path.join(OUT, name)
    fig.savefig(path, dpi=DPI, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    print(f"  ✓ {name}")
    return path

# Chart A — Seats won: 2021 vs 2026 grouped bar
all_p = [p for p in TOP_PARTIES + ['Others']
         if win21.groupby('party_group').size().get(p, 0) > 0
         or win26.groupby('party_group').size().get(p, 0) > 0]
s21 = win21.groupby('party_group').size()
s26 = win26.groupby('party_group').size()

fig, ax = plt.subplots(figsize=(12, 6))
x, w = np.arange(len(all_p)), 0.38
ax.bar(x - w/2, [s21.get(p, 0) for p in all_p], w,
       color=[PARTY_COLORS.get(p, '#999') for p in all_p], alpha=0.45,
       edgecolor='white', linewidth=1.5, label='2021')
b26 = ax.bar(x + w/2, [s26.get(p, 0) for p in all_p], w,
             color=[PARTY_COLORS.get(p, '#999') for p in all_p],
             edgecolor='white', linewidth=1.5, label='2026')
for bar in b26:
    h = bar.get_height()
    if h > 0:
        ax.text(bar.get_x() + bar.get_width()/2, h + 0.8, str(int(h)),
                ha='center', va='bottom', fontsize=10, fontweight='bold')
ax.set_xticks(x)
ax.set_xticklabels(all_p, fontsize=11, fontweight='bold')
ax.set_ylabel('Seats Won', fontsize=12)
ax.set_title('Seats Won by Party — 2021 vs 2026', fontsize=15, fontweight='bold', pad=15, color=NAVY)
ax.legend(['2021', '2026'], fontsize=11, framealpha=0.85)
ax.set_ylim(0, max(s26.max(), s21.max()) * 1.2)
chart_a = save_chart(fig, 'chart_a_party_seats.png')

# Chart B — Regional seats stacked bar (2026)
seats_26 = win26.groupby(['region', 'party_group']).size().reset_index(name='seats')
pivot_b   = seats_26.pivot(index='region', columns='party_group', values='seats').fillna(0)
show_cols = [p for p in TOP_PARTIES + ['Others'] if p in pivot_b.columns]
pivot_b   = pivot_b[show_cols]

fig, ax = plt.subplots(figsize=(13, 6))
bottom = np.zeros(len(pivot_b))
for party in show_cols:
    vals = pivot_b[party].values
    bars = ax.bar(pivot_b.index, vals, bottom=bottom, label=party,
                  color=PARTY_COLORS.get(party, '#999'), edgecolor='white', linewidth=0.8)
    for bar, val, bot in zip(bars, vals, bottom):
        if val >= 3:
            ax.text(bar.get_x() + bar.get_width()/2, bot + val/2,
                    str(int(val)), ha='center', va='center',
                    fontsize=9, color='white', fontweight='bold')
    bottom += vals
ax.set_xlabel('Region', fontsize=12)
ax.set_ylabel('Seats', fontsize=12)
ax.set_title('Seats Won by Party per Region — 2026', fontsize=15, fontweight='bold', pad=15, color=NAVY)
ax.legend(loc='upper right', fontsize=10, framealpha=0.85, ncol=2)
ax.tick_params(axis='x', rotation=15)
chart_b = save_chart(fig, 'chart_b_regional_seats.png')

# Chart C — Flip heatmap
flips      = cmp[cmp['flipped'] == 1]
flip_pivot = flips.groupby(['party_21_group', 'party_26_group']).size().unstack(fill_value=0)

fig, ax = plt.subplots(figsize=(11, 7))
sns.heatmap(flip_pivot, annot=True, fmt='d', cmap='YlOrRd', ax=ax,
            linewidths=0.5, linecolor='white', cbar_kws={'label': 'Seats'})
ax.set_title(f'Seat Flips: 2021 → 2026   ({n_flipped} of {n_const} changed hands)',
             fontsize=14, fontweight='bold', pad=15, color=NAVY)
ax.set_xlabel('Won by in 2026', fontsize=11)
ax.set_ylabel('Held by in 2021', fontsize=11)
chart_c = save_chart(fig, 'chart_c_flip_heatmap.png')

# Chart D — Vote share distribution
fig, ax = plt.subplots(figsize=(11, 6))
ax.hist(win21['vote_share'], bins=20, alpha=0.55, color='#1565C0',
        label=f'2021  (avg {avg21}%)', edgecolor='white')
ax.hist(win26['vote_share'], bins=20, alpha=0.80, color='#FF6B00',
        label=f'2026  (avg {avg26}%)', edgecolor='white')
ax.axvline(35, color='#555', linestyle='--', linewidth=1.8, label='35% threshold')
ax.axvline(50, color='#111', linestyle='--', linewidth=1.8, label='50% majority')
ax.set_xlabel('Winner Vote Share %', fontsize=12)
ax.set_ylabel('Number of Constituencies', fontsize=12)
ax.set_title('Winner Vote Share Distribution — 2021 vs 2026',
             fontsize=15, fontweight='bold', pad=15, color=NAVY)
ax.legend(fontsize=11)
ylim = ax.get_ylim()[1]
ax.text(29, ylim * 0.88, f'Below 35%\n2021: {below35_21}  |  2026: {below35_26}',
        fontsize=9, color='#333',
        bbox=dict(facecolor='#FFF3E0', edgecolor='#FF9800', boxstyle='round,pad=0.35'))
ax.text(50.6, ylim * 0.88, f'Above 50%\n2021: {above50_21}  |  2026: {above50_26}',
        fontsize=9, color='#333',
        bbox=dict(facecolor='#E8F5E9', edgecolor='#4CAF50', boxstyle='round,pad=0.35'))
chart_d = save_chart(fig, 'chart_d_vote_share_dist.png')

# Chart E — Top 15 constituency swings
swing_df = cmp[cmp['flipped'] == 1].copy()
swing_df['abs_swing'] = swing_df['share_swing'].abs()
top_sw = swing_df.nlargest(15, 'abs_swing')[
    ['constituency', 'party_21', 'party_26', 'share_swing']
].sort_values('share_swing')
colors_e = ['#C62828' if s < 0 else '#2E7D32' for s in top_sw['share_swing']]

fig, ax = plt.subplots(figsize=(12, 7))
ax.barh(top_sw['constituency'], top_sw['share_swing'],
        color=colors_e, edgecolor='white', linewidth=0.8)
ax.axvline(0, color='#333', linewidth=1.2)
for (_, row), color in zip(top_sw.iterrows(), colors_e):
    label = f"{row['party_21']} → {row['party_26']}"
    xpos  = row['share_swing']
    y     = top_sw.index.get_loc(row.name)
    # Anchor near the zero line (not the bar tip) so long bars never push the
    # label into the y-axis constituency labels on the far left.
    if xpos >= 0:
        ax.text(-0.3, y, label, va='center', ha='right', fontsize=8, color='#444')
    else:
        ax.text(0.3, y, label, va='center', ha='left', fontsize=8, color='#444')
ax.set_xlabel('Vote Share Swing (2026 − 2021) %', fontsize=11)
ax.set_title('Top 15 Constituency Swings — Flipped Seats',
             fontsize=14, fontweight='bold', pad=15, color=NAVY)
ax.legend(handles=[
    mpatches.Patch(color='#2E7D32', label='2026 winner gained share'),
    mpatches.Patch(color='#C62828', label='2026 winner entered from low base'),
], fontsize=10)
chart_e = save_chart(fig, 'chart_e_top_swings.png')

# Chart F — Avg winner share by region, 2021 vs 2026
reg21 = win21.groupby('region')['vote_share'].mean().round(1)
reg26 = win26.groupby('region')['vote_share'].mean().round(1)
regions = sorted(set(reg21.index) | set(reg26.index))

fig, ax = plt.subplots(figsize=(12, 5.4))
x, w = np.arange(len(regions)), 0.38
ax.bar(x - w/2, [reg21.get(r, 0) for r in regions], w,
       color='#1565C0', alpha=0.65, edgecolor='white', label='2021')
b26r = ax.bar(x + w/2, [reg26.get(r, 0) for r in regions], w,
              color='#FF6B00', edgecolor='white', label='2026')
for bar in b26r:
    ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.4,
            f'{bar.get_height():.1f}%',
            ha='center', va='bottom', fontsize=9, fontweight='bold')
ax.set_xticks(x)
ax.set_xticklabels(regions, fontsize=11)
ax.set_ylabel('Avg Winner Vote Share %', fontsize=11)
ax.set_ylim(0, 70)
ax.set_title('Average Winner Vote Share by Region — 2021 vs 2026',
             fontsize=14, fontweight='bold', pad=15, color=NAVY)
ax.legend(fontsize=11)
chart_f = save_chart(fig, 'chart_f_regional_margins.png')

print("✅ All 6 charts generated")

# ═══════════════════════════════════════════════════════════════════════
# 4. POWERPOINT PRESENTATION
# ═══════════════════════════════════════════════════════════════════════
print("\nBuilding PowerPoint...")

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = 6   # index of blank layout in default Office theme

# ── PPT helper functions ──────────────────────────────────────────────

def rect(slide, left, top, width, height, rgb):
    """Filled colour rectangle (implemented as filled textbox)."""
    sh = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    sh.fill.solid()
    sh.fill.fore_color.rgb = rgb
    sh.line.color.rgb = rgb   # hide border by matching fill colour
    return sh


FONT_BODY = 'Segoe UI'      # clean, modern — body copy, labels, tables
FONT_DISPLAY = 'Cambria'    # serif — cover/section headlines, for an editorial
                             # "official report" register rather than a slide deck

_slide_no = {'n': 0}        # running page counter, incremented by footer()


def txt(slide, text, left, top, width, height,
        size=14, bold=False, italic=False,
        color=RGBColor(0x21, 0x21, 0x21),
        align=PP_ALIGN.LEFT, wrap=True, bg=None, font=FONT_BODY):
    """Add a text box, optionally with a background fill."""
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    if bg:
        tb.fill.solid()
        tb.fill.fore_color.rgb = bg
        tb.line.color.rgb = bg
    else:
        tb.fill.background()
    tf = tb.text_frame
    tf.word_wrap = wrap
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size    = Pt(size)
        run.font.bold    = bold
        run.font.italic  = italic
        run.font.name    = font
        run.font.color.rgb = color
    return tb


def header(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 1.28, C_DARK)
    rect(slide, 0, 1.28, 13.33, 0.04, C_ORANGE)
    txt(slide, title,    0.28, 0.09, 12.6, 0.76, size=25, bold=True,
        color=C_WHITE, font=FONT_DISPLAY)
    if subtitle:
        txt(slide, subtitle, 0.28, 0.83, 12.6, 0.44, size=12, italic=True, color=RGBColor(0xD9, 0xC9, 0x8E))


def footer(slide, note="Tamil Nadu Assembly Elections — Data Analysis 2026", total=10):
    _slide_no['n'] += 1
    rect(slide, 0, 7.25, 13.33, 0.25, RGBColor(0xEC, 0xED, 0xF1))
    txt(slide, note, 0.2, 7.26, 10.5, 0.22, size=9, color=C_GRAY)
    txt(slide, f'Page {_slide_no["n"]} of {total}', 11.4, 7.26, 1.73, 0.22,
        size=9, color=C_GRAY, align=PP_ALIGN.RIGHT)


def pic(slide, path, left, top, width, height=None):
    if height:
        slide.shapes.add_picture(path, Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    else:
        slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width))


def kpi(slide, left, top, value, label, color=C_DARK, bw=3.0, bh=1.65):
    rect(slide, left,        top,      bw,    bh,    C_LIGHT)
    rect(slide, left,        top,      0.07,  bh,    color)
    txt(slide,  str(value),  left+0.15, top+0.1, bw-0.2, 0.82,
        size=30, bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(slide,  label,       left+0.1,  top+0.92, bw-0.2, 0.65,
        size=10, color=C_GRAY, align=PP_ALIGN.CENTER)


# ── Slide 1: Title ────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[BLANK])
rect(s, 0, 0, 13.33, 7.5, C_DARK)
rect(s, 0, 0,    13.33, 0.1,  C_ORANGE)
rect(s, 0, 5.6,  13.33, 0.08, C_ORANGE)
rect(s, 0, 0.1,  0.5,   5.5,  C_ORANGE)
txt(s, 'TAMIL NADU',          0.75, 1.1,  12, 1.1,  size=54, bold=True,
    color=C_WHITE, font=FONT_DISPLAY)
txt(s, 'ASSEMBLY ELECTIONS',  0.75, 2.15, 12, 0.85, size=30, bold=False,
    color=RGBColor(0xC7, 0xCE, 0xDE), font=FONT_DISPLAY)
txt(s, '2  0  2  6',          0.75, 2.95, 12, 1.2,  size=62, bold=True,
    color=C_ORANGE, font=FONT_DISPLAY)
txt(s, 'A Comprehensive Data Analysis', 0.75, 4.2, 11, 0.6,
    size=18, italic=True, color=RGBColor(0xC7, 0xCE, 0xDE))
txt(s, 'Comparing 2021 vs 2026  ·  234 Constituencies  ·  6 Regions  ·  32 Districts',
    0.75, 4.82, 11.5, 0.5, size=12, color=RGBColor(0x9A, 0xA3, 0xBD))
txt(s, 'June 2026', 0.75, 6.75, 5, 0.4, size=11,
    color=RGBColor(0x77, 0x80, 0x9C))

# ── Slide 2: At a Glance ─────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[BLANK])
header(s, 'At a Glance', 'Key numbers from the 2026 Tamil Nadu Assembly Election')
footer(s)

# Row 1 — x positions: 0.28, 3.46, 6.64, 9.82  (width 3.0, gap 0.18)
kpi(s, 0.28, 1.55, '234',       'Total\nConstituencies',    C_DARK)
kpi(s, 3.46, 1.55, str(top_seats), f'Seats Won\n({top_party})', C_ORANGE)
kpi(s, 6.64, 1.55, str(n_flipped), 'Seats\nFlipped',         C_RED)
kpi(s, 9.82, 1.55, f'{avg26}%', 'Avg Winner\nShare 2026',   C_GREEN)

# Row 2
kpi(s, 0.28, 3.38, str(above50_26), 'Winners\n> 50% Share',    C_GREEN)
kpi(s, 3.46, 3.38, str(below35_26), 'Winners\n< 35% Share',    C_RED)
kpi(s, 6.64, 3.38, '6',          'Regions\nCovered',          C_DARK)
kpi(s, 9.82, 3.38, '32',         'Districts\nCovered',        RGBColor(0x4A, 0x14, 0x8C))

txt(s,
    f'★  Average winner vote share fell from {avg21}% (2021) to {avg26}% (2026)'
    f' — elections became significantly more competitive.',
    0.28, 5.22, 12.77, 0.52, size=12, color=C_DARK,
    bg=RGBColor(0xFF, 0xF8, 0xE1))

# ── Slide 3: The 2026 Mandate ─────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[BLANK])
header(s, 'The 2026 Mandate', 'Seats won by party — comparison with 2021')
footer(s)
pic(s, chart_a, 0.15, 1.4, 13.0, 5.65)

# ── Slide 4: Regional Dominance ───────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[BLANK])
header(s, 'Regional Dominance', "Seats won per party across Tamil Nadu's 6 regions — 2026")
footer(s)
pic(s, chart_b, 0.15, 1.4, 13.0, 5.65)

# ── Slide 5: The Flip Story ───────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[BLANK])
header(s, 'The Flip Story',
       f'{n_flipped} of {n_const} constituencies changed hands — who gained, who lost?')
footer(s)
pic(s, chart_c, 0.15, 1.4, 8.2, 5.65)

# Insight panel
rect(s, 8.5, 1.4, 4.68, 5.65, C_LIGHT)
txt(s, 'Key Insights', 8.68, 1.55, 4.3, 0.48, size=14, bold=True, color=C_DARK)
rect(s, 8.5, 2.0, 4.68, 0.05, C_ORANGE)
insights = (
    f'• {n_flipped} seats ({round(n_flipped/n_const*100)}%) changed party\n'
    f'• {n_const - n_flipped} seats retained same party\n'
    '\n'
    '• Row = party that HELD seat in 2021\n'
    '• Column = party that WON in 2026\n'
    '• Number = seats that moved between\n'
    '  those two parties\n'
    '\n'
    '• TVK made significant inroads into\n'
    '  seats previously held by other parties\n'
    '\n'
    '• AIADMK was the biggest loser,\n'
    '  surrendering seats across regions'
)
txt(s, insights, 8.68, 2.1, 4.3, 4.8, size=11, color=C_TEXT)

# ── Slide 6: Margin of Victory ────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[BLANK])
header(s, 'Margin of Victory',
       'Are elections getting more competitive? Vote share distribution 2021 vs 2026')
footer(s)
pic(s, chart_d, 0.15, 1.4, 9.6, 5.7)

# Stats side panel
rect(s, 9.95, 1.4, 3.23, 5.7, C_LIGHT)
txt(s, 'Statistics', 10.12, 1.55, 2.9, 0.45, size=13, bold=True, color=C_DARK)
rect(s, 9.95, 1.98, 3.23, 0.05, C_ORANGE)

stats = [
    ('Avg Share  2021', f'{avg21}%'),
    ('Avg Share  2026', f'{avg26}%'),
    ('', ''),
    ('>50% winners 2021', str(above50_21)),
    ('>50% winners 2026', str(above50_26)),
    ('', ''),
    ('<35% winners 2021', str(below35_21)),
    ('<35% winners 2026', str(below35_26)),
]
for i, (label, val) in enumerate(stats):
    y = 2.1 + i * 0.58
    if label:
        txt(s, label, 10.08, y, 2.0, 0.48, size=10, color=C_GRAY)
        txt(s, val,   12.08, y, 1.0, 0.48, size=11, bold=True,
            color=C_DARK, align=PP_ALIGN.RIGHT)

# ── Slide 7: Top Swings ───────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[BLANK])
header(s, 'Biggest Swings',
       'Constituencies with the largest vote-share change in flipped seats')
footer(s)
pic(s, chart_e, 0.15, 1.4, 13.0, 5.8)

# ── Slide 8: Regional Competitiveness ────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[BLANK])
header(s, 'Regional Competitiveness',
       'Average winner vote share by region — 2021 vs 2026')
footer(s)
pic(s, chart_f, 0.15, 1.4, 13.0, 5.6)

# ── Slide 9: Key Takeaways ────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[BLANK])
header(s, 'Key Takeaways',
       'Summary of findings from the 2026 Tamil Nadu Assembly Election')
footer(s)

takeaways = [
    ('Electoral Shift',
     f'{top_party} emerged as the largest party with {top_seats} seats, '
     f'representing a significant realignment in Tamil Nadu politics.'),
    ('Seat Churn',
     f'{n_flipped} of {n_const} constituencies ({round(n_flipped/n_const*100)}%) '
     f'changed hands — indicating a strong anti-incumbency wave.'),
    ('Closer Contests',
     f'Average winner vote share fell from {avg21}% (2021) to {avg26}% (2026), '
     f'reflecting highly competitive multi-cornered contests.'),
    ('Fragmented Mandate',
     f'{below35_26} winners secured victory with under 35% of the vote — '
     f'a sign of vote splitting across new parties and independents.'),
    ('Regional Patterns',
     'No single party swept all 6 regions; dominance is geographically concentrated, '
     'making regional coalition arithmetic crucial.'),
]

for i, (head, body) in enumerate(takeaways):
    y = 1.46 + i * 1.07
    rect(s, 0.3, y, 0.07, 0.92, C_ORANGE)
    txt(s, head, 0.5, y,       12.6, 0.44, size=13, bold=True, color=C_DARK)
    txt(s, body, 0.5, y + 0.42, 12.6, 0.55, size=11, color=C_TEXT)

# ── Slide 10: Thank You ───────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[BLANK])
rect(s, 0, 0, 13.33, 7.5, C_DARK)
rect(s, 0, 0,   0.5, 7.5, C_ORANGE)
rect(s, 0, 3.55, 13.33, 0.08, C_ORANGE)
txt(s, 'Thank You',            0.75, 2.1,  12, 1.3,  size=58, bold=True,
    color=C_WHITE, font=FONT_DISPLAY)
txt(s, 'Questions  &  Discussion', 0.75, 3.7, 12, 0.8,
    size=20, italic=True, color=RGBColor(0xC7, 0xCE, 0xDE), font=FONT_DISPLAY)
txt(s,
    'Data: Tamil Nadu Election Commission  |  '
    'Analysis: TN Election 2026 Dashboard  |  June 2026',
    0.75, 6.65, 12, 0.45, size=11,
    color=RGBColor(0x77, 0x80, 0x9C))

# ── Save ──────────────────────────────────────────────────────────────
ppt_path = os.path.join(OUT, 'TN_Election_2026_Professional.pptx')
prs.save(ppt_path)
print(f"✅ PowerPoint saved: {ppt_path}")

# ── Summary ───────────────────────────────────────────────────────────
print("\n🎉 ALL DONE — Output files:")
for f in sorted(os.listdir(OUT)):
    size = os.path.getsize(os.path.join(OUT, f)) // 1024
    print(f"  {f:<50}  {size:>6} KB")
